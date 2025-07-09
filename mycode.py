import os
import pandas as pd
import tempfile
from flask import Flask, render_template, request, redirect, flash, jsonify
import google.generativeai as genai
import PyPDF2
import json
import xml.etree.ElementTree as ET
import xml.dom.minidom as minidom
from pydrive.auth import GoogleAuth
from pydrive.drive import GoogleDrive
import gspread
from google.oauth2.service_account import Credentials
from pptx import Presentation

app = Flask(__name__)

# Google Gemini API setup
GEMINI_API_KEY = "AIzaSyB_23f1Z2w6Qun6G52OMg2Gi7QnlFVopGI"
genai.configure(api_key=GEMINI_API_KEY)

model = genai.GenerativeModel('gemini-2.0-flash-exp')
COURSES_EXCEL_FILE = 'Copy of HwKaosi COurse lists.xlsx'

GOOGLE_SHEET_URL = 'https://docs.google.com/spreadsheets/d/1ZFumz6XDdFMauqPu-X-4U5wpLjfytwP_NqdcjaBdYTc/edit?gid=0#gid=0'

PP_XML_FILE_folder = '1tfsGfIrDrVMJhdDyGNlsAeBz7bXFtSNt'  # Your existing PP folder ID
FF_XML_FILE_folder = '1-2ZYcGXlzsCPIFu3W9jTEMFTmILOnp-7'  # Your existing FF folder ID
FINAL_EXAM_XML_FILE_folder = '1-2ZYcGXlzsCPIFu3W9jTEMFTmILOnp-7'  # Replace with your Final Exam folder ID

def extract_text_from_pdf(pdf_file):
    """Extract text from a PDF file"""
    text = ""
    pdf_reader = PyPDF2.PdfReader(pdf_file)
    for page in pdf_reader.pages:
        page_text = page.extract_text()
        if page_text:
            text += page_text + "\n"
    return text


def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)

        for slide_number, slide in enumerate(prs.slides):
            text += f"Slide {slide_number + 1}:\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() != "":
                    text += f"{shape.text.strip()}\n"

            text += "\n" + "-" * 40 + "\n"

        return text
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""


def get_google_sheets_client():
    scope = [
        "https://spreadsheets.google.com/feeds",
        "https://www.googleapis.com/auth/drive"
    ]
    creds = Credentials.from_service_account_file(
        "your-credentials.json",  # Make sure this file exists in your project directory
        scopes=scope
    )
    return gspread.authorize(creds)


def generate_mcqs(text, num_questions=10, custom_instructions=""):
    # Base prompt
    base_prompt = f"""
    Based on the following clinical or educational text, generate {num_questions} NCLEX-style multiple-choice questions.

    Each question must:
    1. Present a realistic clinical scenario or stem.
    2. Provide **four options**: option1 through option4.
    3. Use the `answers` field to list one or more correct options, e.g., ["option2"] or ["option2", "option4"].
    4. Include an `explanation` of the correct answer.
    5. Include a helpful `hint`.

    Format the output strictly as JSON, using 3-space indentation, matching this structure:
    {{
      "title": "An apt title for the exam",
      "questions": [
       {{
        "question": "the question text goes here",
        "option1": "option 1 text",
        "option2": "option 2 text",
        "option3": "option 3 text",
        "option4": "option 4 text",
        "answers": ["option2"],
        "explanation": "The explanation of the answer goes here",
        "hint": "The hint text goes here"
       }}
      ]
    }}
    """

    if custom_instructions and custom_instructions.strip():
        base_prompt += f"\nAdditional instructions for question generation:\n{custom_instructions.strip()}\n"

    full_prompt = base_prompt + f"\nTEXT:\n{text[:8000]}"

    try:
        response = model.generate_content(full_prompt)
        response_text = response.text.strip()
        import re
        json_match = re.search(r'\[\s*{.*}\s*\]', response_text, re.DOTALL)

        if json_match:
            mcq_json = json_match.group(0)
            return json.loads(mcq_json)
        else:
            print("No valid JSON found in the response.")
            return []
    except Exception as e:
        print(f"Error generating MCQs with model {model._model_name}: {e}")
        return []


def json_to_wpadvquiz_xml(nclex_data, quiz_title=None):
    import xml.etree.ElementTree as ET
    from xml.dom import minidom

    # If the input is a list, treat it as just questions
    if isinstance(nclex_data, list):
        questions = nclex_data
        title = quiz_title or "NCLEX Practice Exam"
    else:
        questions = nclex_data.get("questions", [])
        title = quiz_title or nclex_data.get("title", "NCLEX Practice Exam")

    root = ET.Element("wpAdvQuiz")
    header = ET.SubElement(root, "header", version="1.0.1", exportVersion="1")
    data = ET.SubElement(root, "data")
    quiz = ET.SubElement(data, "quiz")

    ET.SubElement(quiz, "title", titleHidden="false").text = f"<![CDATA[{title}]]>"
    questions_node = ET.SubElement(quiz, "questions")

    for idx, mcq in enumerate(questions, 1):
        correct_keys = set(mcq.get("answers", []))
        is_multiple = len(correct_keys) > 1
        question_type = "multiple" if is_multiple else "single"

        question = ET.SubElement(questions_node, "question", answerType=question_type)
        ET.SubElement(question, "title").text = f"<![CDATA[NCLEX Question {idx}]]>"
        ET.SubElement(question, "points").text = "1"
        ET.SubElement(question, "questionText").text = f"<![CDATA[{mcq['question']}]]>"

        answers_node = ET.SubElement(question, "answers")

        for key in ["option1", "option2", "option3", "option4"]:
            option_text = mcq.get(key, "")
            is_correct = "true" if key in correct_keys else "false"
            answer = ET.SubElement(answers_node, "answer", points="1", correct=is_correct)
            ET.SubElement(answer, "answerText", html="false").text = f"<![CDATA[{option_text}]]>"
            ET.SubElement(answer, "stortText", html="false").text = "<![CDATA[]]>"

        if mcq.get("explanation"):
            ET.SubElement(question, "correctMsg").text = f"<![CDATA[{mcq['explanation']}]]>"
        if mcq.get("hint"):
            ET.SubElement(question, "incorrectMsg").text = f"<![CDATA[{mcq['hint']}]]>"

    rough_xml = ET.tostring(root, encoding="utf-8")
    return minidom.parseString(rough_xml).toprettyxml(indent="  ")





def load_courses():
    """Load available courses from Excel file"""
    try:
        if os.path.exists(COURSES_EXCEL_FILE):
            df = pd.read_excel(COURSES_EXCEL_FILE)
            if len(df.columns) > 0:
                return df.iloc[:, 0].tolist()  # First column contains course names
        return []
    except Exception as e:
        print(f"Error loading courses: {e}")
        return []


def upload_to_drive(xml_content, file_name="quiz.xml", quiz_type="foundational"):
    """
    Upload XML content to Google Drive and create a shareable link
    Based on quiz_type:
    - "foundational": upload to FF_XML_FILE_folder
    - "pp": upload to PP_XML_FILE_folder
    - "final_exam": upload to FINAL_EXAM_XML_FILE_folder
    """
    temp_path = None
    try:
        # Create a temporary file with the XML content
        with tempfile.NamedTemporaryFile(delete=False, suffix='.xml') as temp_file:
            temp_path = temp_file.name
            temp_file.write(xml_content.encode('utf-8'))
            temp_file.close()

        gauth = GoogleAuth()
        gauth.LoadCredentialsFile("mycreds.txt")

        if gauth.credentials is None:
            gauth.LocalWebserverAuth()
        elif gauth.access_token_expired:
            gauth.Refresh()
        else:
            gauth.Authorize()

        gauth.SaveCredentialsFile("mycreds.txt")
        drive = GoogleDrive(gauth)

        # Determine which folder to use based on quiz type
        if quiz_type == "pp":
            folder_id = PP_XML_FILE_folder
        elif quiz_type == "final_exam":
            folder_id = FINAL_EXAM_XML_FILE_folder
        else:  # default to foundational
            folder_id = FF_XML_FILE_folder

        file_to_upload = drive.CreateFile({
            'title': file_name,
            'parents': [{'id': folder_id}]
        })

        file_to_upload.SetContentFile(temp_path)
        file_to_upload.Upload()

        # Create a shareable link
        file_to_upload.InsertPermission({
            'type': 'anyone',
            'value': 'anyone',
            'role': 'reader'
        })

        return {
            'success': True,
            'file_id': file_to_upload['id'],
            'view_url': f"https://drive.google.com/file/d/{file_to_upload['id']}/view",
            'share_url': file_to_upload['alternateLink'],
            'quiz_type': quiz_type
        }
    except Exception as e:
        print(f"Error uploading to Google Drive: {e}")
        return {
            'success': False,
            'error': str(e)
        }
    finally:
        if temp_path and os.path.exists(temp_path):
            try:
                os.unlink(temp_path)
            except Exception as e:
                print(f"Warning: Could not delete temporary file {temp_path}: {e}")


def update_sheet_with_status(course_number, share_url=None, quiz_type="foundational"):
    """
    Update the Google Sheet based on quiz type:
    - foundational: store share_url in column P (index 16)
    - pp: store share_url in column S (index 19)
    - final_exam: store share_url in column V (index 22) - adjust this index as needed
    """
    try:
        client = get_google_sheets_client()
        sheet = client.open_by_url(GOOGLE_SHEET_URL)
        worksheet = sheet.sheet1

        # Assuming course numbers are in column A (index 1)
        course_number_column = 1
        all_course_numbers = worksheet.col_values(course_number_column)

        if course_number in all_course_numbers:
            row_index = all_course_numbers.index(course_number) + 1  # 1-based indexing

            if share_url:
                # Determine which column to update based on quiz type
                if quiz_type == "pp":
                    # PP quiz type - update column S (index 19)
                    column_index = 19  # Column S
                    column_letter = "S"
                elif quiz_type == "final_exam":
                    # Final Exam quiz type - update column V (index 22) - adjust as needed
                    column_index = 22  # Column V - adjust this value based on your sheet
                    column_letter = "V"
                else:
                    # Default to Foundational (FF) - update column P (index 16)
                    column_index = 16  # Column P
                    column_letter = "P"

                # Update the appropriate column with the shareable link
                worksheet.update_cell(row_index, column_index, share_url)
                return True, f"Updated column {column_letter} for course number '{course_number}'"

            return False, "No share URL provided"

        return False, f"Course number '{course_number}' not found in sheet"

    except Exception as e:
        print(f"Error updating Google Sheet: {e}")
        return False, f"Error: {str(e)}"


@app.route('/')
def home():
    """Render the home page with tool options"""
    return render_template('home.html')


@app.route('/quiz-generator', methods=['GET', 'POST'])
def quiz_generator():
    """Handle the PDF to Quiz Generator functionality"""
    result = None
    # Load available courses from Excel file
    available_courses = load_courses()

    if request.method == 'POST':
        selected_course = request.form.get('course', '')
        custom_instructions = request.form.get('custom_instructions', '').strip()
        quiz_type = request.form.get('quiz_type', 'foundational')

        if 'file' not in request.files:
            flash('No file part')
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            flash('No selected file')
            return redirect(request.url)

        if file and file.filename.endswith('.pdf'):
            original_filename = os.path.splitext(file.filename)[0]

            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                file.save(temp_file.name)
                pdf_text = extract_text_from_pdf(temp_file.name)

            os.unlink(temp_file.name)

            num_questions = int(request.form.get('num_questions', 10))

            mcqs = generate_mcqs(pdf_text, num_questions, custom_instructions)

            # Use the selected course number as part of the file name
            quiz_file_name = f"{selected_course}_{quiz_type}_quiz.xml" if selected_course else f"{original_filename}_{quiz_type}_quiz.xml"

            xml_content = json_to_wpadvquiz_xml(mcqs, original_filename)
            # Pass quiz_type to upload_to_drive
            drive_result = upload_to_drive(xml_content, quiz_file_name, quiz_type)

            sheet_updated = False
            sheet_message = ""

            if drive_result['success'] and selected_course:
                # Update Google Sheet with share URL in the appropriate column based on quiz type
                success, message = update_sheet_with_status(
                    selected_course,
                    drive_result.get('share_url'),
                    quiz_type
                )

                sheet_updated = success
                quiz_type_display = "Foundational (FF)" if quiz_type == "foundational" else "PP" if quiz_type == "pp" else "Final Exam"
                sheet_message = f" and Google Sheet updated for {quiz_type_display} quiz" if success else f" but Google Sheet update failed: {message}"
            elif not selected_course:
                sheet_message = " but no course selected for Google Sheet update."

            result = {
                'success': drive_result['success'],
                'message': f'Quiz generated and uploaded successfully{sheet_message}' if drive_result[
                    'success'] else f'Error uploading quiz: {drive_result.get("error", "Unknown error")}',
                'view_url': drive_result.get('view_url', ''),
                'sheet_updated': sheet_updated,
                'quiz_type': quiz_type
            }

    return render_template('quiz_generator.html', result=result, courses=available_courses)


import re


@app.route('/generate-pptx-script', methods=['POST'])
def generate_pptx_script():
    file = request.files.get('file')
    custom_prompt = request.form.get('prompt', '')

    if not file or not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Please upload a valid PowerPoint file'}), 400

    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
            file.save(temp_file.name)
            extracted_text = extract_text_from_pptx(temp_file.name)

        os.unlink(temp_file.name)

        if not extracted_text.strip():
            return jsonify({'error': 'No text could be extracted from the PowerPoint file'}), 400

        if "{{extracted_text}}" in custom_prompt:
            final_prompt = custom_prompt.replace("{{extracted_text}}", extracted_text[:8000])
        else:
            final_prompt = f"{custom_prompt}\n\nPowerPoint Content:\n{extracted_text[:8000]}"

        print("Final prompt sent to Gemini:")
        print(final_prompt)

        try:
            ai_response = model.generate_content(final_prompt)
            script = ai_response.text.strip()

            # === Clean the script ===
            clean_lines = []
            for line in script.splitlines():
                stripped_line = line.strip()

                # Skip unwanted markdown/code markers
                if stripped_line.startswith('```') or stripped_line.startswith('##'):
                    continue

                # Remove timestamps like (0:00 - 0:30)
                stripped_line = re.sub(r'\(\s*\d+:\d+\s*-\s*\d+:\d+\s*\)', '', stripped_line)

                # Remove prefixes like 'Narrator:' and 'Visual:'
                stripped_line = re.sub(r'^(Narrator:|Visual:|Script:)\s*', '', stripped_line, flags=re.IGNORECASE)

                # Remove stray asterisks
                stripped_line = stripped_line.replace('*', '')

                if stripped_line:  # Only add non-empty lines
                    clean_lines.append(stripped_line)

            cleaned_script = '\n'.join(clean_lines).strip()

            return jsonify({'script': cleaned_script})
        except Exception as e:
            return jsonify({'error': f'Error generating script with Gemini API: {str(e)}'}), 500

    except Exception as e:
        return jsonify({'error': f'Unexpected server error: {str(e)}'}), 500


@app.route('/', methods=['GET', 'POST'])
def index():
    result = None
    # Load available courses from Excel file
    available_courses = load_courses()

    if request.method == 'POST':
        selected_course = request.form.get('course', '')
        custom_instructions = request.form.get('custom_instructions', '').strip()

        if 'file' not in request.files:
            flash('No file part')
            return redirect(request.url)

        file = request.files['file']
        if file.filename == '':
            flash('No selected file')
            return redirect(request.url)

        if file and file.filename.endswith('.pdf'):
            original_filename = os.path.splitext(file.filename)[0]

            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                file.save(temp_file.name)
                pdf_text = extract_text_from_pdf(temp_file.name)

            os.unlink(temp_file.name)

            num_questions = int(request.form.get('num_questions', 5))

            mcqs = generate_mcqs(pdf_text, num_questions, custom_instructions)

            # Use the selected course name as part of the file name
            quiz_file_name = f"{selected_course}_quiz.xml" if selected_course else f"{original_filename}_quiz.xml"

            xml_content = json_to_wpadvquiz_xml(mcqs, original_filename)
            # Pass custom_instructions to upload_to_drive
            drive_result = upload_to_drive(xml_content, quiz_file_name, custom_instructions)

            excel_updated = False
            if drive_result['success'] and selected_course:
                # Update spreadsheet with the share URL in appropriate column
                excel_updated, message = update_sheet_with_status(
                    selected_course,
                    drive_result['share_url'],
                    custom_instructions
                )

                column_letter = "S" if custom_instructions.strip() else "P"
                excel_message = f" and Excel file column {column_letter} updated." if excel_updated else " but Excel update failed."
            else:
                excel_message = ""
                excel_updated = False

            result = {
                'success': drive_result['success'],
                'message': f'Quiz generated and uploaded successfully{excel_message}' if drive_result[
                    'success'] else 'Error uploading quiz.',
                'view_url': drive_result.get('view_url', ''),
                'excel_updated': excel_updated
            }

    return render_template('index.html', result=result, courses=available_courses)


def extract_text_from_pptx(file_path):
    text = ""
    try:
        prs = Presentation(file_path)

        for slide_number, slide in enumerate(prs.slides):
            text += f"Slide {slide_number + 1}:\n"

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip() != "":
                    text += f"{shape.text.strip()}\n"

            text += "\n" + "-" * 40 + "\n"

        return text
    except Exception as e:
        print(f"Error extracting text from PowerPoint: {e}")
        return ""


@app.route('/generate-pdf-script', methods=['POST'])
def generate_pdf_script():
    file = request.files.get('file')
    custom_prompt = request.form.get('prompt', '')

    if not file or not file.filename.lower().endswith('.pdf'):
        return jsonify({'error': 'Please upload a valid PDF file'}), 400

    try:
        # Save PDF temporarily and extract text
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            file.save(temp_file.name)
            extracted_text = extract_text_from_pdf(temp_file.name)

        os.unlink(temp_file.name)

        if not extracted_text.strip():
            return jsonify({'error': 'No text could be extracted from the PDF'}), 400

        # Replace placeholder in custom prompt with extracted text
        final_prompt = custom_prompt.replace("{{extracted_text}}", extracted_text[:8000])

        # Generate script using Gemini API
        try:
            ai_response = model.generate_content(final_prompt)
            script = ai_response.text.strip()
            return jsonify({'script': script})
        except Exception as e:
            return jsonify({'error': f'Error generating script with Gemini API: {str(e)}'}), 500

    except Exception as e:
        return jsonify({'error': f'Error processing PDF: {str(e)}'}), 500


@app.route('/pdf-to-video')
def pdf_to_video():
    return render_template('pptx_to_video.html')  # Or rename the template to 'pdf_to_video.html'


os.makedirs('templates', exist_ok=True)

with open('templates/home.html', 'w', encoding="utf-8") as f:
    f.write("""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>EduTools Hub</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700&display=swap');

        :root {
            --primary: #6C63FF;
            --primary-light: #837CFF;
            --secondary: #2D3142;
            --light: #F5F5F5;
            --grey: #9BA4B4;
            --dark: #21232F;
            --success: #00C896;
            --error: #FF6B6B;
            --warning: #FFB648;
        }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        body {
            font-family: 'Poppins', sans-serif;
            background: linear-gradient(135deg, var(--dark) 0%, #1A1C29 100%);
            color: var(--light);
            line-height: 1.6;
            min-height: 100vh;
            display: flex;
            flex-direction: column;
            align-items: center;
            padding: 40px 20px;
        }

        .header {
            text-align: center;
            margin-bottom: 50px;
            animation: fadeInDown 0.8s ease-out;
        }

        @keyframes fadeInDown {
            from { opacity: 0; transform: translateY(-30px); }
            to { opacity: 1; transform: translateY(0); }
        }

        .header-logo {
            font-size: 42px;
            font-weight: 700;
            color: var(--primary);
            margin-bottom: 15px;
            display: flex;
            align-items: center;
            justify-content: center;
        }

        .header-logo i {
            margin-right: 15px;
        }

        h1 {
            font-size: 32px;
            font-weight: 600;
            margin-bottom: 10px;
            color: white;
        }

        .subtitle {
            color: var(--grey);
            font-size: 18px;
            font-weight: 400;
            max-width: 600px;
            margin: 0 auto;
        }

        .tools-container {
            display: flex;
            justify-content: center;
            gap: 30px;
            margin-top: 30px;
            flex-wrap: wrap;
            max-width: 1200px;
            width: 100%;
        }

        .tool-card {
            background: rgba(45, 49, 66, 0.8);
            backdrop-filter: blur(10px);
            border-radius: 16px;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
            padding: 30px;
            width: 100%;
            max-width: 400px;
            border: 1px solid rgba(255, 255, 255, 0.1);
            transition: all 0.3s ease;
            position: relative;
            overflow: hidden;
            animation: fadeIn 0.6s ease-out;
            animation-fill-mode: backwards;
        }

        .tool-card:nth-child(1) {
            animation-delay: 0.2s;
        }

        .tool-card:nth-child(2) {
            animation-delay: 0.4s;
        }

        @keyframes fadeIn {
            from { opacity: 0; transform: translateY(20px); }
            to { opacity: 1; transform: translateY(0); }
        }

        .tool-card:hover {
            transform: translateY(-10px);
            box-shadow: 0 15px 35px rgba(0, 0, 0, 0.2);
            border-color: rgba(108, 99, 255, 0.3);
        }

        .tool-card::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            width: 100%;
            height: 4px;
            background: linear-gradient(90deg, var(--primary) 0%, var(--primary-light) 100%);
        }

        .tool-icon {
            font-size: 48px;
            color: var(--primary);
            margin-bottom: 20px;
            text-align: center;
        }

        .tool-title {
            font-size: 22px;
            font-weight: 600;
            margin-bottom: 12px;
            text-align: center;
        }

        .tool-description {
            color: var(--grey);
            font-size: 15px;
            margin-bottom: 25px;
            text-align: center;
            min-height: 75px;
        }

        .tool-features {
            margin-bottom: 25px;
        }

        .feature-item {
            display: flex;
            align-items: center;
            margin-bottom: 10px;
        }

        .feature-item i {
            color: var(--success);
            margin-right: 10px;
            font-size: 14px;
        }

        .tool-button {
            display: block;
            width: 100%;
            padding: 14px;
            background: var(--primary);
            color: white;
            border: none;
            border-radius: 8px;
            font-size: 16px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            text-align: center;
            text-decoration: none;
        }

        .tool-button:hover {
            background: var(--primary-light);
            transform: translateY(-2px);
            box-shadow: 0 6px 12px rgba(108, 99, 255, 0.2);
        }

        .tool-button:active {
            transform: translateY(0);
        }

        .coming-soon {
            background: rgba(255, 182, 72, 0.2);
            color: var(--warning);
            cursor: not-allowed;
        }

        .coming-soon:hover {
            background: rgba(255, 182, 72, 0.3);
            transform: none;
            box-shadow: none;
        }

        .footer {
            margin-top: 60px;
            text-align: center;
            color: var(--grey);
            font-size: 14px;
        }

        @media (max-width: 768px) {
            .tools-container {
                flex-direction: column;
                align-items: center;
            }

            .tool-card {
                max-width: 100%;
            }

            .header {
                margin-bottom: 30px;
            }

            h1 {
                font-size: 28px;
            }

            .subtitle {
                font-size: 16px;
            }
        }
    </style>
</head>
<body>
    <div class="header">
        <div class="header-logo">
            <i class="fas fa-graduation-cap"></i>
            EduTools Hub
        </div>
        <h1>AI-Powered Educational Tools</h1>
        <p class="subtitle">Transform your educational content with our intelligent conversion tools designed for educators and trainers</p>
    </div>

    <div class="tools-container">
        <!-- PDF to Quiz Tool Card -->
        <div class="tool-card">
            <div class="tool-icon">
                <i class="fas fa-brain"></i>
            </div>
            <div class="tool-title">PDF to Quiz Generator</div>
            <div class="tool-description">
                Convert your PDF documents into interactive quizzes with multiple-choice questions using AI technology.
            </div>
            <div class="tool-features">
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>AI-generated multiple-choice questions</span>
                </div>
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>Customizable question settings</span>
                </div>
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>WordPress Quiz XML format export</span>
                </div>
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>Google Drive integration</span>
                </div>
            </div>
            <a href="/quiz-generator" class="tool-button">Create Quiz Now</a>
        </div>

        <!-- PPTX to Video Tool Card -->
        <div class="tool-card">
            <div class="tool-icon">
                <i class="fas fa-film"></i>
            </div>
            <div class="tool-title">PPTX to Video Generator</div>
            <div class="tool-description">
                Transform your PowerPoint presentations into engaging video content with voiceovers and animations.
            </div>
            <div class="tool-features">
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>Automatic slide animations</span>
                </div>
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>AI-generated voiceover narration</span>
                </div>
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>Multiple output formats</span>
                </div>
                <div class="feature-item">
                    <i class="fas fa-check-circle"></i>
                    <span>Custom timing and transitions</span>
                </div>
            </div>
            <a href="/pptx-to-video" class="tool-button">Convert Now</a>
        </div>
    </div>

    <div class="footer">
        <p>EduTools Hub • AI-Powered Educational Tools • © 2025</p>
    </div>
</body>
</html>
    """)

with open('templates/index.html', 'w', encoding="utf-8") as f:
    f.write("""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Smart Quiz Generator</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <style>
        @import url('https://fonts.googleapis.com/css2?family=Poppins:wght@300;400;500;600;700&display=swap');

        :root {
            --primary: #6C63FF;
            --primary-light: #837CFF;
            --secondary: #2D3142;
            --light: #F5F5F5;
            --grey: #9BA4B4;
            --dark: #21232F;
            --success: #00C896;
            --error: #FF6B6B;
            --warning: #FFB648;
        }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }

        body {
            font-family: 'Poppins', sans-serif;
            background: linear-gradient(135deg, var(--dark) 0%, #1A1C29 100%);
            color: var(--light);
            line-height: 1.6;
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            padding: 20px;
        }

        .container {
            background: rgba(45, 49, 66, 0.8);
            backdrop-filter: blur(10px);
            border-radius: 16px;
            box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
            padding: 40px;
            max-width: 800px;
            width: 100%;
            border: 1px solid rgba(255, 255, 255, 0.1);
            animation: fadeIn 0.6s ease-out;
        }

        @keyframes fadeIn {
            from { opacity: 0; transform: translateY(20px); }
            to { opacity: 1; transform: translateY(0); }
        }

        .header {
            text-align: center;
            margin-bottom: 30px;
        }

        .logo {
            font-size: 32px;
            font-weight: 700;
            color: var(--primary);
            margin-bottom: 8px;
            display: flex;
            align-items: center;
            justify-content: center;
        }

        .logo i {
            margin-right: 12px;
        }

        h1 {
            font-size: 24px;
            font-weight: 600;
            margin-bottom: 8px;
            color: white;
        }

        .subtitle {
            color: var(--grey);
            font-size: 16px;
            font-weight: 400;
        }

        .form-group {
            margin-bottom: 20px;
        }

        label {
            display: block;
            margin-bottom: 8px;
            font-weight: 500;
            font-size: 15px;
        }

        .input-container {
            position: relative;
            border-radius: 8px;
            overflow: hidden;
            transition: all 0.3s ease;
            border: 1px solid rgba(255, 255, 255, 0.1);
        }

        .input-container:focus-within {
            border-color: var(--primary);
            box-shadow: 0 0 0 3px rgba(108, 99, 255, 0.2);
        }

        .file-input-container {
            position: relative;
            background: rgba(255, 255, 255, 0.05);
            border-radius: 8px;
            padding: 15px;
            text-align: center;
            cursor: pointer;
            border: 2px dashed rgba(108, 99, 255, 0.3);
            transition: all 0.3s ease;
        }

        .file-input-container:hover {
            border-color: var(--primary);
            background: rgba(108, 99, 255, 0.05);
        }

        .file-input-container i {
            font-size: 36px;
            color: var(--primary);
            margin-bottom: 10px;
        }

        .file-input-text {
            font-weight: 500;
            margin-bottom: 5px;
        }

        .file-input-subtext {
            font-size: 13px;
            color: var(--grey);
        }

        input[type="file"] {
            position: absolute;
            width: 100%;
            height: 100%;
            top: 0;
            left: 0;
            opacity: 0;
            cursor: pointer;
        }

        select {
            width: 100%;
            padding: 12px 16px;
            font-size: 15px;
            border: none;
            background-color: rgba(255, 255, 255, 0.05);
            color: var(--light);
            border-radius: 8px;
            cursor: pointer;
            appearance: none;
            -webkit-appearance: none;
            transition: all 0.3s ease;
        }

        select:focus {
            outline: none;
            background-color: rgba(255, 255, 255, 0.1);
        }

        textarea {
            width: 100%;
            min-height: 100px;
            padding: 12px 16px;
            font-size: 15px;
            border: none;
            background-color: rgba(255, 255, 255, 0.05);
            color: var(--light);
            border-radius: 8px;
            resize: vertical;
            font-family: 'Poppins', sans-serif;
            transition: all 0.3s ease;
        }

        textarea:focus {
            outline: none;
            background-color: rgba(255, 255, 255, 0.1);
            border-color: var(--primary);
        }

        .select-container {
            position: relative;
        }

        .select-container::after {
            content: '\f078';
            font-family: 'Font Awesome 6 Free';
            font-weight: 900;
            position: absolute;
            right: 16px;
            top: 50%;
            transform: translateY(-50%);
            color: var(--grey);
            pointer-events: none;
        }

        button {
            width: 100%;
            padding: 14px;
            background: var(--primary);
            color: white;
            border: none;
            border-radius: 8px;
            font-size: 16px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            display: flex;
            justify-content: center;
            align-items: center;
        }

        button i {
            margin-right: 8px;
        }

        button:hover {
            background: var(--primary-light);
            transform: translateY(-2px);
            box-shadow: 0 6px 12px rgba(108, 99, 255, 0.2);
        }

        button:active {
            transform: translateY(0);
        }

        .features {
            margin-top: 40px;
            display: grid;
            grid-template-columns: repeat(3, 1fr);
            gap: 20px;
        }

        .feature {
            text-align: center;
            padding: 20px;
            background: rgba(255, 255, 255, 0.03);
            border-radius: 8px;
            transition: all 0.3s ease;
        }

        .feature:hover {
            transform: translateY(-5px);
            background: rgba(255, 255, 255, 0.05);
        }

        .feature i {
            font-size: 28px;
            color: var(--primary);
            margin-bottom: 12px;
        }

        .feature-title {
            font-weight: 600;
            margin-bottom: 5px;
        }

        .feature-desc {
            font-size: 13px;
            color: var(--grey);
        }

        .footer {
            margin-top: 40px;
            text-align: center;
            color: var(--grey);
            font-size: 14px;
        }

        .instruction-help {
            color: var(--grey);
            font-size: 13px;
            margin-top: 5px;
            font-style: italic;
        }

        .result-container {
            background: rgba(0, 200, 150, 0.1);
            border-left: 4px solid var(--success);
            border-radius: 6px;
            padding: 16px;
            margin-bottom: 25px;
            animation: fadeIn 0.5s ease;
        }

        .result-container.error {
            background: rgba(255, 107, 107, 0.1);
            border-left: 4px solid var(--error);
        }

        .result-title {
            font-weight: 600;
            margin-bottom: 8px;
            display: flex;
            align-items: center;
        }

        .result-title i {
            margin-right: 8px;
        }

        .result-message {
            margin-bottom: 12px;
        }

        .result-link {
            display: inline-flex;
            align-items: center;
            background: rgba(255, 255, 255, 0.1);
            padding: 8px 16px;
            border-radius: 4px;
            color: var(--light);
            text-decoration: none;
            font-size: 14px;
            font-weight: 500;
            transition: all 0.2s ease;
        }

        .result-link:hover {
            background: rgba(255, 255, 255, 0.15);
        }

        .result-link i {
            margin-right: 8px;
        }

        @media (max-width: 768px) {
            .container {
                padding: 30px 20px;
            }

            .features {
                grid-template-columns: 1fr;
            }
        }
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <div class="logo">
                <i class="fas fa-brain"></i>
                Smart Quiz
            </div>
            <h1>PDF to Quiz Generator</h1>
            <p class="subtitle">Upload a PDF and generate quiz questions instantly</p>
        </div>

        {% if result %}
        <div class="result-container {% if not result.success %}error{% endif %}">
            <div class="result-title">
                {% if result.success %}
                <i class="fas fa-check-circle"></i> Success!
                {% else %}
                <i class="fas fa-exclamation-circle"></i> Error
                {% endif %}
            </div>
            <div class="result-message">{{ result.message }}</div>
            {% if result.view_url %}
            <a href="{{ result.view_url }}" target="_blank" class="result-link">
                <i class="fas fa-external-link-alt"></i> View Quiz on Google Drive
            </a>
            {% endif %}
        </div>
        {% endif %}

        <form method="POST" enctype="multipart/form-data">
            <div class="form-group">
                <label for="file">Upload PDF Document</label>
                <div class="file-input-container">
                    <i class="fas fa-file-pdf"></i>
                    <div class="file-input-text">Drag & drop your PDF here</div>
                    <div class="file-input-subtext">or click to browse files</div>
                    <input type="file" name="file" id="file" accept=".pdf" required>
                </div>
            </div>

            <div class="form-group">
                <label for="course">Select Course</label>
                <div class="select-container">
                    <select name="course" id="course">
                        <option value="">-- Select a Course --</option>
                        {% for course in courses %}
                        <option value="{{ course }}">{{ course }}</option>
                        {% endfor %}
                    </select>
                </div>
            </div>
            
            <div class="form-group">
                <label for="quiz_type">Quiz Type</label>
                <div class="select-container">
                    <select name="quiz_type" id="quiz_type">
                        <option value="foundational">Foundational (FF)</option>
                        <option value="pp">PP</option>
                        <option value="final_exam">Final Exam</option>
                    </select>
                </div>
            </div>
            
            <div class="form-group">
                <label for="num_questions">Number of Questions</label>
                <div class="select-container">
                    <select name="num_questions" id="num_questions">
                        <option value="10">10 Questions</option>
                        <option value="20">20 Questions</option>
                        <option value="30">30 Questions</option>
                        <option value="40">40 Questions</option>
                    </select>
                </div>
            </div>

            <div class="form-group">
                <label for="custom_instructions">Custom MCQ Instructions</label>
                <textarea name="custom_instructions" id="custom_instructions" placeholder="Example: Make the questions challenging, focus on key concepts, include calculation problems, etc."></textarea>
                <p class="instruction-help">Provide specific instructions on how you want your quiz questions to be generated</p>
            </div>

            <button type="submit" id="generateBtn">
                <i class="fas fa-magic"></i>
                Generate Quiz
            </button>
        </form>

        <div class="features">
            <div class="feature">
                <i class="fas fa-bolt"></i>
                <div class="feature-title">Instant Generation</div>
                <div class="feature-desc">Generate quizzes within seconds</div>
            </div>
            <div class="feature">
                <i class="fas fa-bullseye"></i>
                <div class="feature-title">Accurate Questions</div>
                <div class="feature-desc">AI-powered relevant MCQs</div>
            </div>
            <div class="feature">
                <i class="fas fa-cloud-upload-alt"></i>
                <div class="feature-title">Google Drive Export</div>
                <div class="feature-desc">Quiz saved to your Drive</div>
            </div>
        </div>

        <div class="footer">
            <p>Powered by AI • PDF Quiz Generator • © 2025</p>
        </div>
    </div>

    <script>
        document.querySelector('form').addEventListener('submit', function() {
            const button = document.getElementById('generateBtn');
            button.innerHTML = '<i class="fas fa-spinner fa-spin"></i> Generating...';
            button.disabled = true;
        });
    </script>
</body>
</html>
    """)

with open('templates/quiz_generator.html', 'w', encoding="utf-8") as f:
    f.write("""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PDF to Quiz Generator - EduTools Hub</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <style>
        :root {
            --primary-blue: #2E7BFF;
            --dark-blue: #1a1e2e;
            --accent-blue: #56CCF2;
            --background-dark: #0F1223;
            --text-light: #f8f9fa;
            --text-dim: #b0b8c2;
            --success-color: #10b981;
            --error-color: #ef4444;
        }
        
        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }
        
        body {
            background: linear-gradient(135deg, var(--background-dark) 0%, #111827 100%);
            color: var(--text-light);
            min-height: 100vh;
            display: flex;
            flex-direction: column;
            line-height: 1.6;
            padding: 20px;
        }
        
        .container {
            width: 90%;
            max-width: 800px;
            padding: 40px;
            background: rgba(26, 32, 53, 0.7);
            backdrop-filter: blur(10px);
            border-radius: 16px;
            box-shadow: 0 20px 50px rgba(0, 0, 0, 0.5), 
                        0 0 0 1px rgba(255, 255, 255, 0.05),
                        0 0 30px rgba(46, 123, 255, 0.2);
            position: relative;
            overflow: hidden;
            margin: 20px auto;
            animation: fadeIn 0.6s ease-out;
        }
        
        @keyframes fadeIn {
            from { opacity: 0; transform: translateY(20px); }
            to { opacity: 1; transform: translateY(0); }
        }
        
        .container::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: linear-gradient(90deg, var(--primary-blue), var(--accent-blue));
        }
        
        .nav-back {
            color: var(--text-light);
            text-decoration: none;
            display: flex;
            align-items: center;
            font-size: 15px;
            transition: all 0.2s ease;
            margin-bottom: 10px;
            width: fit-content;
            margin-left: 5%;
        }
        
        .nav-back:hover {
            color: var(--accent-blue);
            transform: translateX(-3px);
        }
        
        .nav-back i {
            margin-right: 8px;
        }
        
        .header {
            text-align: center;
            margin-bottom: 30px;
            position: relative;
        }
        
        .logo {
            font-size: 32px;
            font-weight: 700;
            margin-bottom: 8px;
            display: flex;
            align-items: center;
            justify-content: center;
            background: linear-gradient(90deg, var(--primary-blue), var(--accent-blue));
            -webkit-background-clip: text;
            background-clip: text;
            -webkit-text-fill-color: transparent;
        }
        
        .logo i {
            margin-right: 12px;
            -webkit-text-fill-color: var(--accent-blue);
        }
        
        h1 {
            color: white;
            text-align: center;
            margin-bottom: 10px;
            font-weight: 700;
            font-size: 2.2rem;
            letter-spacing: -0.5px;
        }
        
        .subtitle {
            color: var(--text-dim);
            font-size: 1.1rem;
        }
        
        .form-group {
            margin-bottom: 24px;
        }
        
        label {
            display: block;
            margin-bottom: 10px;
            font-weight: 500;
            font-size: 16px;
            color: var(--accent-blue);
        }
        
        .file-input-container {
            position: relative;
            background: rgba(26, 30, 46, 0.6);
            border-radius: 12px;
            padding: 25px;
            text-align: center;
            cursor: pointer;
            border: 2px dashed rgba(86, 204, 242, 0.3);
            transition: all 0.3s ease;
        }
        
        .file-input-container:hover {
            border-color: var(--primary-blue);
            background: rgba(46, 123, 255, 0.1);
        }
        
        .file-input-container i {
            font-size: 40px;
            color: var(--accent-blue);
            margin-bottom: 12px;
        }
        
        .file-input-text {
            font-weight: 500;
            margin-bottom: 6px;
            font-size: 17px;
        }
        
        .file-input-subtext {
            font-size: 14px;
            color: var(--text-dim);
        }
        
        input[type="file"] {
            position: absolute;
            width: 100%;
            height: 100%;
            top: 0;
            left: 0;
            opacity: 0;
            cursor: pointer;
        }

        input[type="number"] {
            width: 100%;
            padding: 14px 16px;
            font-size: 16px;
            border: 1px solid rgba(86, 204, 242, 0.2);
            background-color: rgba(26, 30, 46, 0.6);
            color: var(--text-light);
            border-radius: 12px;
            transition: all 0.3s ease;
        }
        
        input[type="number"]:focus {
            outline: none;
            border-color: var(--primary-blue);
            box-shadow: 0 0 0 3px rgba(46, 123, 255, 0.2);
        }
        
        .select-container {
            position: relative;
        }
        
        .select-container::after {
            content: '\f078';
            font-family: 'Font Awesome 6 Free';
            font-weight: 900;
            position: absolute;
            right: 16px;
            top: 50%;
            transform: translateY(-50%);
            color: var(--text-dim);
            pointer-events: none;
            font-size: 14px;
        }
        
        select {
            width: 100%;
            padding: 14px 16px;
            font-size: 16px;
            border: none;
            background-color: rgba(26, 30, 46, 0.6);
            color: var(--text-light);
            border-radius: 12px;
            cursor: pointer;
            appearance: none;
            -webkit-appearance: none;
            transition: all 0.3s ease;
            border: 1px solid rgba(86, 204, 242, 0.2);
        }
        
        select:focus {
            outline: none;
            border-color: var(--primary-blue);
            box-shadow: 0 0 0 3px rgba(46, 123, 255, 0.2);
        }
        
        textarea {
            width: 100%;
            min-height: 120px;
            padding: 14px 16px;
            font-size: 16px;
            border: 1px solid rgba(86, 204, 242, 0.2);
            background-color: rgba(26, 30, 46, 0.6);
            color: var(--text-light);
            border-radius: 12px;
            resize: vertical;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
            transition: all 0.3s ease;
        }
        
        textarea:focus {
            outline: none;
            border-color: var(--primary-blue);
            box-shadow: 0 0 0 3px rgba(46, 123, 255, 0.2);
        }
        
        .instruction-help {
            color: var(--text-dim);
            font-size: 14px;
            margin-top: 8px;
            font-style: italic;
        }
        
        button {
            width: 100%;
            padding: 16px;
            background: linear-gradient(135deg, var(--primary-blue), var(--accent-blue));
            color: white;
            border: none;
            border-radius: 12px;
            font-size: 16px;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            display: flex;
            justify-content: center;
            align-items: center;
            box-shadow: 0 4px 20px rgba(46, 123, 255, 0.4);
        }
        
        button i {
            margin-right: 10px;
        }
        
        button:hover {
            transform: translateY(-2px);
            box-shadow: 0 8px 25px rgba(46, 123, 255, 0.5);
        }
        
        button:active {
            transform: translateY(1px);
        }
        
        button:disabled {
            background: linear-gradient(135deg, #4b5563, #6b7280);
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }
        
        .result-container {
            background: rgba(16, 185, 129, 0.15);
            border-left: 4px solid var(--success-color);
            border-radius: 12px;
            padding: 20px;
            margin-bottom: 30px;
            animation: fadeIn 0.5s ease;
        }
        
        .result-container.error {
            background: rgba(239, 68, 68, 0.15);
            border-left: 4px solid var(--error-color);
        }
        
        .result-title {
            font-weight: 600;
            margin-bottom: 10px;
            display: flex;
            align-items: center;
            font-size: 18px;
        }
        
        .result-title i {
            margin-right: 10px;
            font-size: 20px;
        }
        
        .result-message {
            margin-bottom: 15px;
            color: var(--text-dim);
        }
        
        .result-link {
            display: inline-flex;
            align-items: center;
            background: rgba(255, 255, 255, 0.1);
            padding: 10px 18px;
            border-radius: 8px;
            color: var(--text-light);
            text-decoration: none;
            font-size: 15px;
            font-weight: 500;
            transition: all 0.3s ease;
            border: 1px solid rgba(255, 255, 255, 0.1);
        }
        
        .result-link:hover {
            background: rgba(46, 123, 255, 0.2);
            transform: translateY(-2px);
        }
        
        .result-link i {
            margin-right: 10px;
        }
        
        .features {
            margin-top: 40px;
            display: grid;
            grid-template-columns: repeat(3, 1fr);
            gap: 20px;
        }
        
        .feature {
            text-align: center;
            padding: 24px 15px;
            background: rgba(26, 30, 46, 0.6);
            border-radius: 14px;
            transition: all 0.3s ease;
            border: 1px solid rgba(86, 204, 242, 0.1);
        }
        
        .feature:hover {
            transform: translateY(-5px);
            border-color: rgba(86, 204, 242, 0.3);
            box-shadow: 0 10px 20px rgba(0, 0, 0, 0.2);
        }
        
        .feature i {
            font-size: 32px;
            color: var(--accent-blue);
            margin-bottom: 15px;
        }
        
        .feature-title {
            font-weight: 600;
            margin-bottom: 8px;
            font-size: 17px;
        }
        
        .feature-desc {
            font-size: 14px;
            color: var(--text-dim);
        }
        
        .footer {
            margin-top: 40px;
            text-align: center;
            color: var(--text-dim);
            font-size: 14px;
        }
        
        /* Glowing effect */
        .glow {
            position: absolute;
            width: 300px;
            height: 300px;
            border-radius: 50%;
            background: radial-gradient(circle, rgba(46, 123, 255, 0.2) 0%, rgba(0, 0, 0, 0) 70%);
            z-index: -1;
            filter: blur(30px);
        }
        
        .glow-1 {
            top: -150px;
            left: -150px;
        }
        
        .glow-2 {
            bottom: -100px;
            right: -70px;
            background: radial-gradient(circle, rgba(86, 204, 242, 0.15) 0%, rgba(0, 0, 0, 0) 70%);
        }
        
        @media (max-width: 768px) {
            .container {
                padding: 30px 20px;
            }
            
            .features {
                grid-template-columns: 1fr;
            }
            
            h1 {
                font-size: 1.8rem;
            }
            
            .logo {
                font-size: 28px;
            }
        }
    </style>
</head>
<body>
    <a href="/" class="nav-back"><i class="fas fa-arrow-left"></i> Back to Home</a>

    <div class="container">
        <div class="glow glow-1"></div>
        <div class="glow glow-2"></div>
        
        <div class="header">
            <div class="logo">
                <i class="fas fa-brain"></i>
                Smart Quiz
            </div>
            <h1>PDF to Quiz Generator</h1>
            <p class="subtitle">Upload a PDF and generate quiz questions instantly</p>
        </div>

        {% if result %}
        <div class="result-container {% if not result.success %}error{% endif %}">
            <div class="result-title">
                {% if result.success %}
                <i class="fas fa-check-circle"></i> Success!
                {% else %}
                <i class="fas fa-exclamation-circle"></i> Error
                {% endif %}
            </div>
            <div class="result-message">{{ result.message }}</div>
            {% if result.view_url %}
            <a href="{{ result.view_url }}" target="_blank" class="result-link">
                <i class="fas fa-external-link-alt"></i> View Quiz on Google Drive
            </a>
            {% endif %}
        </div>
        {% endif %}

        <form method="POST" enctype="multipart/form-data">
            <div class="form-group">
                <label for="file">Upload PDF Document</label>
                <div class="file-input-container">
                    <i class="fas fa-file-pdf"></i>
                    <div class="file-input-text">Drag & drop your PDF here</div>
                    <div class="file-input-subtext">or click to browse files</div>
                    <input type="file" name="file" id="file" accept=".pdf" required>
                </div>
            </div>

            <div class="form-group">
                <label for="course">Select Course</label>
                <div class="select-container">
                    <select name="course" id="course">
                        <option value="">-- Select a Course --</option>
                        {% for course in courses %}
                        <option value="{{ course }}">{{ course }}</option>
                        {% endfor %}
                    </select>
                </div>
            </div>

            <!-- New Quiz Type Selection -->
            <div class="form-group">
                <label for="quiz_type">Quiz Type</label>
                <div class="select-container">
                    <select name="quiz_type" id="quiz_type">
                        <option value="foundational">Foundational (FF)</option>
                        <option value="pp">PP</option>
                        <option value="final_exam">Final Exam</option>
                    </select>
                </div>
            </div>

            <div class="form-group">
                <label for="num_questions">Number of Questions</label>
                <input type="number" name="num_questions" id="num_questions" min="1" max="100" value="10" placeholder="Enter number of questions">
            </div>

            <div class="form-group">
                <label for="custom_instructions">Custom MCQ Instructions</label>
                <textarea name="custom_instructions" id="custom_instructions" placeholder="Example: Make the questions challenging, focus on key concepts, include calculation problems, etc."></textarea>
                <p class="instruction-help">Provide specific instructions on how you want your quiz questions to be generated</p>
            </div>

            <button type="submit" id="generateBtn">
                <i class="fas fa-magic"></i>
                Generate Quiz
            </button>
        </form>

        <div class="features">
            <div class="feature">
                <i class="fas fa-bolt"></i>
                <div class="feature-title">Instant Generation</div>
                <div class="feature-desc">Generate quizzes within seconds</div>
            </div>
            <div class="feature">
                <i class="fas fa-bullseye"></i>
                <div class="feature-title">Accurate Questions</div>
                <div class="feature-desc">AI-powered relevant MCQs</div>
            </div>
            <div class="feature">
                <i class="fas fa-cloud-upload-alt"></i>
                <div class="feature-title">Google Drive Export</div>
                <div class="feature-desc">Quiz saved to your Drive</div>
            </div>
        </div>

        <div class="footer">
            <p>Powered by AI • PDF Quiz Generator • © 2025</p>
        </div>
    </div>

    <script>
        document.querySelector('form').addEventListener('submit', function() {
            const button = document.getElementById('generateBtn');
            button.innerHTML = '<i class="fas fa-spinner fa-spin"></i> Generating...';
            button.disabled = true;
        });
    </script>
</body>
</html>
    """)


@app.route('/pptx-to-video')
def pptx_to_video():
    return render_template('pptx_to_video.html')


with open('templates/pptx_to_video.html', 'w', encoding="utf-8") as f:
    f.write("""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>PowerPoint to Educational Script - EduTools Hub</title>
    <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
    <style>
        :root {
            --primary-blue: #2E7BFF;
            --dark-blue: #1a1e2e;
            --accent-blue: #56CCF2;
            --background-dark: #0F1223;
            --text-light: #f8f9fa;
            --text-dim: #b0b8c2;
            --success-color: #10b981;
            --error-color: #ef4444;
        }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
            font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
        }

        body {
            background: linear-gradient(135deg, var(--background-dark) 0%, #111827 100%);
            color: var(--text-light);
            min-height: 100vh;
            display: flex;
            justify-content: center;
            align-items: center;
            line-height: 1.6;
            padding: 20px 0;
        }

        .container {
            width: 95%;
            max-width: 900px;
            padding: 40px;
            background: rgba(26, 32, 53, 0.7);
            backdrop-filter: blur(10px);
            border-radius: 16px;
            box-shadow: 0 20px 50px rgba(0, 0, 0, 0.5), 
                        0 0 0 1px rgba(255, 255, 255, 0.05),
                        0 0 30px rgba(46, 123, 255, 0.2);
            position: relative;
            overflow: hidden;
        }

        .container::before {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            height: 4px;
            background: linear-gradient(90deg, var(--primary-blue), var(--accent-blue));
        }

        h1 {
            color: white;
            text-align: center;
            margin-bottom: 10px;
            font-weight: 700;
            font-size: 2.4rem;
            letter-spacing: -0.5px;
            background: linear-gradient(90deg, var(--accent-blue), #a5f3fc);
            -webkit-background-clip: text;
            background-clip: text;
            -webkit-text-fill-color: transparent;
        }

        p {
            color: var(--text-dim);
            text-align: center;
            margin-bottom: 30px;
            font-size: 1.1rem;
        }

        .file-upload {
            background: rgba(26, 30, 46, 0.6);
            padding: 30px;
            border-radius: 14px;
            text-align: center;
            border: 2px dashed rgba(86, 204, 242, 0.4);
            transition: all 0.3s ease;
            position: relative;
            margin-bottom: 25px;
        }

        .file-upload.highlight {
            border-color: var(--accent-blue);
            background: rgba(46, 123, 255, 0.1);
        }

        .file-upload input {
            width: 0.1px;
            height: 0.1px;
            opacity: 0;
            overflow: hidden;
            position: absolute;
            z-index: -1;
        }

        .file-upload label {
            cursor: pointer;
            display: flex;
            flex-direction: column;
            align-items: center;
            justify-content: center;
            padding: 20px;
        }

        .file-upload i {
            font-size: 3rem;
            color: var(--accent-blue);
            margin-bottom: 15px;
        }

        .file-upload h3 {
            margin-bottom: 10px;
            color: var(--text-light);
        }

        .file-upload p {
            margin-bottom: 0;
            font-size: 0.95rem;
        }

        .upload-info {
            display: none;
            margin-top: 15px;
            padding: 10px;
            background: rgba(16, 185, 129, 0.1);
            border-radius: 8px;
            color: var(--success-color);
        }

        button {
            background: linear-gradient(90deg, var(--primary-blue), var(--accent-blue));
            color: white;
            border: none;
            border-radius: 8px;
            padding: 12px 24px;
            font-size: 1rem;
            font-weight: 600;
            cursor: pointer;
            transition: all 0.3s ease;
            display: block;
            width: 100%;
            margin-bottom: 20px;
        }

        button:hover {
            transform: translateY(-2px);
            box-shadow: 0 5px 15px rgba(46, 123, 255, 0.4);
        }

        button:disabled {
            background: #566175;
            cursor: not-allowed;
            transform: none;
            box-shadow: none;
        }

        #loader {
            display: none;
            text-align: center;
            margin: 20px 0;
        }

        #prompt-container {
            background: rgba(26, 30, 46, 0.6);
            padding: 25px;
            border-radius: 14px;
            margin-top: 20px;
            margin-bottom: 25px;
            text-align: left;
            border: 1px solid rgba(86, 204, 242, 0.2);
        }

        #prompt-container h3 {
            margin-top: 0;
            color: var(--accent-blue);
            font-size: 1.3rem;
            display: flex;
            align-items: center;
            gap: 8px;
            margin-bottom: 15px;
        }

        #prompt-instructions {
            color: var(--text-dim);
            font-size: 0.95rem;
            margin-bottom: 15px;
        }

        #prompt-editor {
            width: 100%;
            min-height: 200px;
            background: rgba(15, 18, 35, 0.7);
            color: var(--text-light);
            border: 1px solid rgba(86, 204, 242, 0.2);
            border-radius: 8px;
            padding: 15px;
            font-family: 'Courier New', monospace;
            font-size: 0.95rem;
            line-height: 1.5;
            resize: vertical;
        }

        .spinner {
            width: 40px;
            height: 40px;
            margin: 0 auto;
            border: 4px solid rgba(255, 255, 255, 0.1);
            border-radius: 50%;
            border-top-color: var(--accent-blue);
            animation: spin 1s linear infinite;
        }

        @keyframes spin {
            to { transform: rotate(360deg); }
        }

        #output {
            display: none;
            background: rgba(26, 30, 46, 0.6);
            padding: 25px;
            border-radius: 14px;
            margin-top: 20px;
            text-align: left;
            border: 1px solid rgba(86, 204, 242, 0.2);
        }

        #output h3 {
            margin-top: 0;
            color: var(--accent-blue);
            font-size: 1.3rem;
            display: flex;
            align-items: center;
            gap: 8px;
            margin-bottom: 15px;
        }

        #scriptContent {
            background: rgba(15, 18, 35, 0.7);
            color: var(--text-light);
            padding: 20px;
            border-radius: 8px;
            white-space: pre-wrap;
            max-height: 400px;
            overflow-y: auto;
            margin-bottom: 15px;
            font-size: 0.95rem;
            line-height: 1.7;
        }

        .copy-btn {
            background: rgba(86, 204, 242, 0.2);
            color: var(--accent-blue);
            border: 1px solid rgba(86, 204, 242, 0.3);
            padding: 8px 16px;
            border-radius: 6px;
            cursor: pointer;
            transition: all 0.2s ease;
            display: flex;
            align-items: center;
            gap: 8px;
            font-size: 0.9rem;
            margin: 0 auto;
        }

        .copy-btn:hover {
            background: rgba(86, 204, 242, 0.3);
        }

        #instructions {
            background: rgba(26, 30, 46, 0.6);
            padding: 25px;
            border-radius: 14px;
            margin-top: 20px;
            text-align: left;
            border: 1px solid rgba(86, 204, 242, 0.2);
            position: relative;
        }

        #instructions h3 {
            margin-top: 0;
            color: var(--accent-blue);
            font-size: 1.3rem;
            display: flex;
            align-items: center;
            gap: 8px;
            margin-bottom: 15px;
        }

        #instructions ol {
            padding-left: 25px;
        }

        #instructions li {
            margin-bottom: 12px;
            color: var(--text-dim);
        }

        #instructions li:last-child {
            margin-bottom: 0;
        }

        .glow {
            position: absolute;
            width: 300px;
            height: 300px;
            border-radius: 50%;
            background: radial-gradient(circle, rgba(46, 123, 255, 0.2) 0%, rgba(0, 0, 0, 0) 70%);
            z-index: -1;
            filter: blur(30px);
        }

        .glow-1 {
            top: -150px;
            left: -150px;
        }

        .glow-2 {
            bottom: -100px;
            right: -70px;
            background: radial-gradient(circle, rgba(86, 204, 242, 0.15) 0%, rgba(0, 0, 0, 0) 70%);
        }

        #error {
            display: none;
            background: rgba(239, 68, 68, 0.1);
            color: var(--error-color);
            padding: 15px;
            border-radius: 8px;
            margin-bottom: 20px;
            text-align: center;
        }

        @media (max-width: 768px) {
            .container {
                padding: 25px;
                width: 95%;
            }

            h1 {
                font-size: 2rem;
            }
        }
    </style>
</head>
<body>
    <div class="container">
        <div class="glow glow-1"></div>
        <div class="glow glow-2"></div>

        <h1>PowerPoint to Educational Script</h1>
        <p>Upload your PowerPoint presentation to generate an educational script for Pictory with Gemini AI.</p>

        <div id="error"></div>

        <div class="file-upload" id="drop-area">
            <input type="file" id="fileInput" accept=".pptx" />
            <label for="fileInput">
                <i class="fas fa-file-powerpoint"></i>
                <h3>Drop your PowerPoint file</h3>
                <p>or click to browse (.pptx)</p>
            </label>
            <div class="upload-info" id="fileInfo"></div>
        </div>

        <div id="prompt-container">
            <h3><i class="fas fa-magic"></i> Customize AI Prompt</h3>
            <p id="prompt-instructions">Edit the prompt below to customize how Gemini AI generates your script. This determines the style, tone accessions for accessibility, and structure of your educational video script.</p>
            <textarea id="prompt-editor">You are an expert educational content creator. Create a well-structured, engaging script for a Pictory text-to-video lecture based on the following PowerPoint content. The script should:
                1. Start with an engaging introduction that outlines what the audience will learn
                2. Break down the content into clear sections with smooth transitions
                3. Use simple, conversational language suitable for students
                4. Incorporate analogies or examples to explain complex concepts
                5. Conclude with a summary of key points and a call to action (e.g., "Try this yourself!" or "Join us next time!")
                6. Now, remember i have to provide this stuff to pictory's ai "text to video generation", make it in just text for the video
                7. Don't write stuff like Narrator, Visual(going to slide 3 etc), also don't write [SCENE CHANGE], just provide  what narrater needs to say and the main topic of it.
                8. Don't write slide number, just provide me data as a script.

                PowerPoint Content:
                {{extracted_text}}
            </textarea>
        </div>

        <button id="generateBtn" disabled>Generate Educational Script</button>

        <div id="loader">
            <div class="spinner"></div>
            <p>Processing your PowerPoint with Gemini AI...</p>
        </div>

        <div id="output">
            <h3><i class="fas fa-magic"></i> Generated Educational Script</h3>
            <div id="scriptContent"></div>
            <button class="copy-btn" id="copyBtn">
                <i class="fas fa-copy"></i> Copy to Clipboard
            </button>
        </div>

        <div id="instructions">
            <h3><i class="fas fa-info-circle"></i> How to Use</h3>
            <ol>
                <li>Upload your PowerPoint presentation using the form above.</li>
                <li>Customize the AI prompt if needed to adjust the style and content of your script.</li>
                <li>Click "Generate Educational Script" to create a student-friendly script with Gemini AI.</li>
                <li>Copy the generated script to your clipboard.</li>
                <li>Paste it into Pictory's text-to-video tool to create your educational video.</li>
            </ol>
        </div>
    </div>

    <script>
        const dropArea = document.getElementById('drop-area');
        const fileInput = document.getElementById('fileInput');
        const fileInfo = document.getElementById('fileInfo');
        const generateBtn = document.getElementById('generateBtn');
        const loader = document.getElementById('loader');
        const output = document.getElementById('output');
        const scriptContent = document.getElementById('scriptContent');
        const copyBtn = document.getElementById('copyBtn');
        const errorDiv = document.getElementById('error');
        const promptEditor = document.getElementById('prompt-editor');

        // Enable button only when a valid file is selected
        fileInput.addEventListener('change', function() {
            handleFiles(this.files);
        });

        function handleFiles(files) {
            if (files.length > 0) {
                const file = files[0];
                if (file.name.toLowerCase().endsWith('.pptx')) {
                    fileInfo.style.display = 'block';
                    fileInfo.innerHTML = `<i class="fas fa-check-circle"></i> ${file.name} selected`;
                    generateBtn.disabled = false;
                    output.style.display = 'none';
                    errorDiv.style.display = 'none';
                } else {
                    showError('Please upload a PowerPoint file (.pptx)');
                    fileInput.value = '';
                    generateBtn.disabled = true;
                }
            }
        }

        // Drag and drop handlers
        ['dragenter', 'dragover', 'dragleave', 'drop'].forEach(eventName => {
            dropArea.addEventListener(eventName, preventDefaults, false);
        });

        function preventDefaults(e) {
            e.preventDefault();
            e.stopPropagation();
        }

        ['dragenter', 'dragover'].forEach(eventName => {
            dropArea.addEventListener(eventName, () => dropArea.classList.add('highlight'), false);
        });

        ['dragleave', 'drop'].forEach(eventName => {
            dropArea.addEventListener(eventName, () => dropArea.classList.remove('highlight'), false);
        });

        dropArea.addEventListener('drop', (e) => {
            const files = e.dataTransfer.files;
            handleFiles(files);
        });

        generateBtn.addEventListener('click', async () => {
            if (!fileInput.files || !fileInput.files.length) {
                showError('Please select a PowerPoint file first.');
                return;
            }

            const file = fileInput.files[0];
            showLoader(true);

            const formData = new FormData();
            formData.append('file', file);
            formData.append('prompt', promptEditor.value);

            try {
                const response = await fetch('/generate-pptx-script', {
                    method: 'POST',
                    body: formData
                });

                const result = await response.json();

                if (result.error) {
                    showError(result.error);
                } else {
                    displayScript(result.script);
                }
            } catch (error) {
                showError('Error processing your PowerPoint: ' + error.message);
            } finally {
                showLoader(false);
            }
        });

        function displayScript(script) {
            scriptContent.textContent = script;
            output.style.display = 'block';
            output.scrollIntoView({ behavior: 'smooth' });
        }

        copyBtn.addEventListener('click', () => {
            const textToCopy = scriptContent.textContent;
            navigator.clipboard.writeText(textToCopy).then(() => {
                const originalText = copyBtn.innerHTML;
                copyBtn.innerHTML = '<i class="fas fa-check"></i> Copied!';
                setTimeout(() => {
                    copyBtn.innerHTML = originalText;
                }, 2000);
            });
        });

        function showLoader(show) {
            loader.style.display = show ? 'block' : 'none';
            generateBtn.disabled = show;
        }

        function showError(message) {
            errorDiv.style.display = 'block';
            errorDiv.innerHTML = `<i class="fas fa-exclamation-circle"></i> ${message}`;
            loader.style.display = 'none';
        }
    </script>
</body>
</html>
    """)

if __name__ == '__main__':
    app.run(debug=True)